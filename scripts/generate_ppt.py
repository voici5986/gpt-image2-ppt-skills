#!/usr/bin/env python3
"""
PPT Generator - Generate PPT slide images using OpenAI gpt-image-2 (Images API).

Generates 16:9 slide images from a slide plan + style template, saves
structured slide_spec metadata for precise editing, and packages a .pptx.
"""

import argparse
import json
import os
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

# =============================================================================
# Constants
# =============================================================================

OUTPUT_BASE_DIR = "outputs"

SCRIPT_DIR = Path(__file__).parent
SKILL_ROOT = SCRIPT_DIR.parent
CWD = Path.cwd()


# =============================================================================
# Environment Configuration
# =============================================================================

SKILL_PREFIX = "gpt-image2-ppt"

SKILL_ENV_MAP = {
    "OPENAI_API_KEY":        "OPENAI_API_KEY",
    "OPENAI_BASE_URL":       "OPENAI_BASE_URL",
    "GPT_IMAGE_MODEL_NAME":  "GPT_IMAGE_MODEL_NAME",
    "GPT_IMAGE_QUALITY":     "GPT_IMAGE_QUALITY",
    "GPT_IMAGE_BACKEND":     "GPT_IMAGE_BACKEND",
    "GPT_IMAGE_CONCURRENCY": "GPT_IMAGE_CONCURRENCY",
    "VISION_BASE_URL":       "VISION_BASE_URL",
    "VISION_API_KEY":        "VISION_API_KEY",
    "VISION_MODEL_NAME":     "VISION_MODEL_NAME",
}


def _load_platform_env() -> bool:
    """Map platform-injected 'skill-prefix_VAR' env vars to unprefixed names.

    The platform (MedAgent / Codex / OpenClaw) sets env vars like
        gpt-image2-ppt_OPENAI_API_KEY=sk-...
    but the script expects plain OPENAI_API_KEY.  This function detects
    the prefixed form and maps them transparently into os.environ.
    Returns True if any prefixed var was found.
    """
    found = False
    for base_name in SKILL_ENV_MAP:
        prefixed = f"{SKILL_PREFIX}_{base_name}"
        value = os.environ.get(prefixed)
        if value and not os.environ.get(base_name):
            os.environ[base_name] = value
            found = True
    return found


def _load_scoped_env_files() -> List[str]:
    """Load env files from explicit or skill-owned locations only.

    This intentionally does not read CWD/.env or walk parent directories. Skill
    credentials should come from the agent framework / system environment first;
    .env is only a standalone CLI fallback.
    """
    candidates: List[Path] = []
    explicit = os.environ.get("GPT_IMAGE2_PPT_ENV")
    if explicit:
        candidates.append(Path(explicit).expanduser())
    candidates.extend([
        SKILL_ROOT / ".env",
        Path.home() / ".codex/skills/gpt-image2-ppt-skills/.env",
        Path.home() / ".claude/skills/gpt-image2-ppt-skills/.env",
        Path.home() / "skills/gpt-image2-ppt/.env",
        Path.home() / "skills/gpt-image2-ppt-skills/.env",
    ])

    try:
        from dotenv import load_dotenv
    except ImportError:
        return []

    loaded: List[str] = []
    seen = set()
    for path in candidates:
        resolved = str(path.resolve()) if path.exists() else str(path)
        if resolved in seen or not path.is_file():
            continue
        seen.add(resolved)
        if load_dotenv(path, override=False):
            loaded.append(resolved)
    return loaded


def load_skill_env() -> None:
    """Load environment with safe precedence.

    Existing process environment wins. Platform-prefixed variables are mapped
    next, then explicit / skill-owned .env files fill in missing values.
    """
    _load_platform_env()
    _load_scoped_env_files()

# =============================================================================
# Session & Metadata Management
# =============================================================================

METADATA_FILENAME = "metadata.json"


def _find_sessions(base_dir: Optional[str] = None) -> List[Dict[str, Any]]:
    """Discover all generation sessions (directories containing metadata.json).

    Returns a list of dicts with 'timestamp', 'dir', 'title', 'slides'.
    """
    if base_dir is None:
        base_dir = str(CWD / OUTPUT_BASE_DIR)
    sessions: List[Dict[str, Any]] = []
    if not os.path.isdir(base_dir):
        return sessions
    for entry in sorted(os.listdir(base_dir), reverse=True):
        session_dir = os.path.join(base_dir, entry)
        meta_path = os.path.join(session_dir, METADATA_FILENAME)
        if os.path.isfile(meta_path):
            try:
                with open(meta_path, "r", encoding="utf-8") as f:
                    meta = json.load(f)
                sessions.append({
                    "timestamp": entry,
                    "dir": session_dir,
                    "title": meta.get("title", "Untitled"),
                    "slide_count": len(meta.get("slides", {})),
                })
            except Exception:
                sessions.append({
                    "timestamp": entry,
                    "dir": session_dir,
                    "title": "(corrupt metadata)",
                    "slide_count": 0,
                })
    return sessions


def _resolve_session(session_id: str) -> str:
    """Resolve a session timestamp or path to the session directory.

    Accepts: full path, relative path, or timestamp in OUTPUT_BASE_DIR.
    """
    # Absolute or relative path
    if os.path.isdir(session_id) and os.path.isfile(os.path.join(session_id, METADATA_FILENAME)):
        return os.path.abspath(session_id)
    # Try under OUTPUT_BASE_DIR
    candidate = str(CWD / OUTPUT_BASE_DIR / session_id)
    if os.path.isdir(candidate) and os.path.isfile(os.path.join(candidate, METADATA_FILENAME)):
        return candidate
    # Try with partial match (prefix)
    sessions = _find_sessions()
    for s in sessions:
        if s["timestamp"].startswith(session_id):
            return s["dir"]
    print(f"[X] Session not found: {session_id}")
    print(f"    Tried: {session_id}")
    if candidate != session_id:
        print(f"    Tried: {candidate}")
    sys.exit(1)


def _load_metadata(session_dir: str) -> Dict[str, Any]:
    """Load metadata.json from a session directory."""
    meta_path = os.path.join(session_dir, METADATA_FILENAME)
    if not os.path.isfile(meta_path):
        print(f"[X] No metadata.json found in {session_dir}")
        sys.exit(1)
    with open(meta_path, "r", encoding="utf-8") as f:
        return json.load(f)


def _save_metadata(metadata: Dict[str, Any], session_dir: str) -> None:
    """Atomically save metadata.json to a session directory.

    Writes to a temp file first, then renames to avoid corruption
    on crashes or disk-full conditions.
    """
    import tempfile as _tempfile
    meta_path = os.path.join(session_dir, METADATA_FILENAME)
    fd, tmp_path = _tempfile.mkstemp(
        suffix=".json", prefix=".metadata_", dir=session_dir
    )
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)
        os.replace(tmp_path, meta_path)  # atomic on POSIX + modern Windows
    except Exception:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass
        raise


def _collect_slide_numbers(metadata: Dict[str, Any]) -> List[int]:
    """Extract sorted slide numbers from metadata slides dict."""
    nums: List[int] = []
    for k in metadata.get("slides", {}).keys():
        try:
            nums.append(int(k))
        except ValueError:
            pass
    nums.sort()
    return nums


def _get_latest_slide_spec(metadata: Dict[str, Any], slide_number: int) -> Dict[str, Any]:
    """Get the current slide_spec for a given slide number."""
    slide_key = str(slide_number)
    slide_data = metadata.get("slides", {}).get(slide_key)
    if not slide_data:
        return {}
    current_version = slide_data.get("current_version", 1)
    for v in slide_data.get("versions", []):
        if v.get("version") == current_version:
            return v.get("spec", {})
    return {}


def _get_version_info(metadata: Dict[str, Any], slide_number: int, version: int) -> Optional[Dict[str, Any]]:
    """Get a specific version entry for a slide."""
    slide_key = str(slide_number)
    slide_data = metadata.get("slides", {}).get(slide_key)
    if not slide_data:
        return None
    for v in slide_data.get("versions", []):
        if v.get("version") == version:
            return v
    return None


def _init_slide_metadata(
    slide_number: int,
    page_type: str,
    initial_spec: Dict[str, Any],
    prompt_file: str,
    image_path: str,
) -> Dict[str, Any]:
    """Create the initial metadata entry for a slide."""
    return {
        "slide_number": slide_number,
        "page_type": page_type,
        "current_version": 1,
        "image_snapshot": image_path,
        "versions": [
            {
                "version": 1,
                "action": "generate",
                "spec": initial_spec,
                "prompt_file": prompt_file,
                "image_snapshot": image_path,
            }
        ],
    }


def _add_slide_version(
    slide_data: Dict[str, Any],
    new_version: int,
    spec: Dict[str, Any],
    action: str,
    image_snapshot: str,
    prompt_file: str,
    edit_instruction: str = "",
    reference_version: Optional[int] = None,
) -> None:
    """Append a new version entry to slide metadata."""
    entry: Dict[str, Any] = {
        "version": new_version,
        "action": action,
        "spec": spec,
        "prompt_file": prompt_file,
        "image_snapshot": image_snapshot,
    }
    if edit_instruction:
        entry["edit_instruction"] = edit_instruction
    if reference_version is not None:
        entry["reference_version"] = reference_version
    slide_data.setdefault("versions", []).append(entry)
    slide_data["current_version"] = new_version
    slide_data["image_snapshot"] = image_snapshot


def _stabilize_version_snapshots(slide_data: Dict[str, Any], slide_number: int, images_dir: str) -> None:
    """Point version history entries at immutable versioned image files when present."""
    for version_entry in slide_data.get("versions", []):
        version = version_entry.get("version")
        if not isinstance(version, int):
            continue
        versioned_rel = f"images/slide-{slide_number:02d}_v{version:04d}.png"
        versioned_abs = os.path.join(images_dir, f"slide-{slide_number:02d}_v{version:04d}.png")
        if os.path.isfile(versioned_abs):
            version_entry["image_snapshot"] = versioned_rel


# =============================================================================
# Slide Spec Helpers
# =============================================================================

def apply_spec_updates(spec: Dict[str, Any], element_updates: Dict[str, Any]) -> Dict[str, Any]:
    """Apply element-level updates to a slide_spec, returning the updated spec.

    element_updates maps element_id to a dict of key-value changes.
    e.g. {"subtitle": {"content": "新内容"}, "title": {"color": "#ff0000"}}
    """
    import copy
    updated = copy.deepcopy(spec)
    elements = updated.setdefault("elements", {})
    for elem_id, changes in element_updates.items():
        if elem_id in elements:
            elements[elem_id].update(changes)
        else:
            elements[elem_id] = changes
    return updated


def construct_edit_prompt(spec: Dict[str, Any], element_updates: Dict[str, Any]) -> str:
    """Construct a natural-language edit prompt from the old spec and element changes.

    The prompt instructs gpt-image-2 to change only specific elements
    while keeping everything else identical to the reference image.
    """
    parts = ["在参考图基础上，只修改以下内容，保持其他所有元素（背景、装饰、布局、字体、颜色、大小）完全不变："]
    elements = spec.get("elements", {})

    for elem_id, changes in element_updates.items():
        old_elem = elements.get(elem_id, {})
        position = old_elem.get("position", "相应位置")
        etype = old_elem.get("type", "元素")
        old_content = old_elem.get("content", "")
        new_content = changes.get("content", old_content)

        if "content" in changes and old_content and new_content != old_content:
            parts.append(
                f"将{position}{etype}的文字从「{old_content}」改为「{new_content}」"
            )
        elif "color" in changes:
            old_color = old_elem.get("color", "")
            new_color = changes.get("color", "")
            parts.append(
                f"将{position}{etype}的颜色从{old_color}改为{new_color}"
            )
        elif "style" in changes:
            parts.append(
                f"将{position}{etype}的样式改为{changes['style']}"
            )
        else:
            for k, v in changes.items():
                parts.append(f"修改{position}{etype}的{k}为{v}")

    return "\n".join(parts)


def generate_prompt_from_spec(
    style_template: str,
    slide_spec: Dict[str, Any],
    page_type: str,
    slide_number: int,
    total_slides: int,
) -> str:
    """Generate a detailed prompt from a structured slide_spec.

    Uses element-level descriptions (type, content, position, style, color)
    to construct a precise prompt that gpt-image-2 can follow.
    """
    elements = slide_spec.get("elements", {})
    layout = slide_spec.get("layout", "")

    # Build element descriptions
    element_lines = []
    for elem_id, elem in elements.items():
        etype = elem.get("type", "unknown")
        content = elem.get("content", "")
        position = elem.get("position", "")
        style_hint = elem.get("style", "")
        color = elem.get("color", "")
        description = elem.get("description", "")

        if description:
            element_lines.append(f"- {elem_id}（{etype}）: {description}")
            continue

        desc_parts = [f"- {etype}「{content}」"]
        if position:
            desc_parts.append(f"位置: {position}")
        if style_hint:
            desc_parts.append(f"样式: {style_hint}")
        if color:
            desc_parts.append(f"颜色: {color}")
        element_lines.append("，".join(desc_parts))

    elements_text = "\n".join(element_lines)

    # Page type hint
    is_cover = page_type == "cover" or slide_number == 1
    is_data = page_type == "data" or slide_number == total_slides
    if is_cover:
        label = "封面页（cover）"
        hint = "标题/副标题处理为视觉焦点，按本风格的封面构图规范处理。"
    elif is_data:
        label = "数据页（data）"
        hint = "突出关键数字、对比或结论；按本风格的数据/总结构图规范处理。"
    else:
        label = "内容页（content）"
        hint = "把要点按本风格的内容构图规范结构化呈现，注意层级、对齐、留白。"

    layout_line = f"\n页面布局: {layout}" if layout else ""

    return (
        style_template
        + "\n\n---\n\n"
        + f"现在请生成本组中的【{label}】，{hint}{layout_line}"
        + "\n\n本页各元素的精确描述（请严格按以下布局、位置、样式生成）：\n\n"
        + elements_text
        + LANGUAGE_FONT_RULE
    )


# =============================================================================
# Style Template
# =============================================================================

def load_style_template(style_path: str) -> str:
    """Extract the '## 基础提示词模板' section from a style markdown file."""
    with open(style_path, "r", encoding="utf-8") as f:
        content = f.read()

    base_prompt_marker = "## 基础提示词模板"
    start_idx = content.find(base_prompt_marker)

    if start_idx == -1:
        print("Warning: '## 基础提示词模板' section not found, using fallback extraction")
        start_idx = content.find("## ")
        end_idx = content.find("## ", start_idx + 3)
        if start_idx == -1 or end_idx == -1:
            return content
        return content[start_idx + 3:end_idx].strip()

    section_start = start_idx + len(base_prompt_marker)
    next_section_idx = content.find("\n## ", section_start)

    if next_section_idx == -1:
        extracted = content[section_start:]
    else:
        extracted = content[section_start:next_section_idx]

    return extracted.strip()


# =============================================================================
# Prompt Generation
# =============================================================================

LANGUAGE_FONT_RULE = """

【强制语言与字体要求】
1. 幻灯片上所有文字必须使用简体中文，严禁出现任何英文单词或句子（产品名称等专有名词可保留英文，其余一律用中文）。
2. 中文字体使用思源黑体（Source Han Sans）或苹方（PingFang SC），字形清晰、笔画规整，严禁使用草书、艺术字或变形字体。
3. 标题字体粗体，正文字体常规，字号对比清晰，确保在演示场景下可读性极高。
"""


def generate_prompt(
    style_template: str,
    page_type: str,
    content_text: str,
    slide_number: int,
    total_slides: int,
    slide_spec: Optional[Dict[str, Any]] = None,
) -> str:
    """Generate a complete prompt for a single slide.

    When slide_spec is provided (the Agent constructed it from plan + style),
    it produces a precise element-by-element prompt.  Otherwise falls back to
    the original freeform content-based prompt.

    Per-page composition rules live inside each style's `## 基础提示词模板`
    (cover / content / data sub-blocks).
    """
    if slide_spec and slide_spec.get("elements"):
        return generate_prompt_from_spec(
            style_template, slide_spec, page_type, slide_number, total_slides
        )

    is_cover = page_type == "cover" or slide_number == 1
    is_data = page_type == "data" or slide_number == total_slides
    if is_cover:
        label = "封面页（cover）"
        hint = "标题/副标题处理为视觉焦点，按本风格的封面构图规范处理。"
    elif is_data:
        label = "数据页（data）"
        hint = "突出关键数字、对比或结论；按本风格的数据/总结构图规范处理。"
    else:
        label = "内容页（content）"
        hint = "把要点按本风格的内容构图规范结构化呈现，注意层级、对齐、留白。"

    return (
        style_template
        + "\n\n---\n\n"
        + f"现在请生成本组中的【{label}】，{hint}\n"
        + "本页要呈现的内容如下（请按本风格美学重新设计版式，不要原样照搬文本节奏）：\n\n"
        + content_text
        + LANGUAGE_FONT_RULE
    )


# =============================================================================
# Image Generation
# =============================================================================

def generate_slide(
    prompt: str,
    slide_number: int,
    output_dir: str,
    reference_image_path: Optional[str] = None,
    backend: str = "openai",
) -> Optional[str]:
    """Generate a single PPT slide image using gpt-image-2.

    backend:
      "openai" (default) -- direct /v1/images or /v1/chat calls, needs OPENAI_API_KEY
      "codex"            -- shell out to `codex exec`, reuses codex CLI auth
    """
    sys.path.insert(0, str(SCRIPT_DIR))
    if backend == "codex":
        from codex_backend import CodexImageBackend as _Backend
    else:
        from image_generator import GptImage2Generator as _Backend

    print(f"  Generating slide {slide_number} via {backend} backend ...")

    generator = _Backend(aspect_ratio="16:9")
    image_path = os.path.join(output_dir, "images", f"slide-{slide_number:02d}.png")

    scene_data = {
        "index": slide_number,
        "image_prompt": prompt,
    }
    generator.generate_scene_image(
        scene_data=scene_data,
        output_path=image_path,
        reference_image_path=reference_image_path,
    )
    print(f"  Slide {slide_number} saved: {image_path}")
    return image_path


# =============================================================================
# Output Generation
# =============================================================================

def save_prompts(output_dir: str, prompts_data: Dict[str, Any]) -> str:
    """Save all prompts to JSON file."""
    prompts_path = os.path.join(output_dir, "prompts.json")
    with open(prompts_path, "w", encoding="utf-8") as f:
        json.dump(prompts_data, f, ensure_ascii=False, indent=2)
    print(f"  Prompts saved: {prompts_path}")
    return prompts_path


def generate_pptx(
    output_dir: str,
    slide_numbers: List[int],
    title: str = "Untitled",
) -> Optional[str]:
    """把 images/slide-XX.png 打包成 16:9 .pptx，每页填满。

    需要 python-pptx；如果没装就跳过并提示。
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
    except ImportError:
        print("(!)  跳过 .pptx 生成（缺 python-pptx，pip install python-pptx 后重试）")
        return None

    prs = Presentation()
    # 标准 16:9 PPT 尺寸：13.333 x 7.5 英寸（1280x720pt）
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # 完全空白布局

    img_dir = os.path.join(output_dir, "images")
    added = 0
    for i in slide_numbers:
        img_path = os.path.join(img_dir, f"slide-{i:02d}.png")
        if not os.path.exists(img_path):
            print(f"  跳过 slide-{i:02d}.png（文件不存在）")
            continue
        slide = prs.slides.add_slide(blank)
        # 图片填满整页（如果原图比例不是 16:9，python-pptx 默认按指定 width/height 拉伸）
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        added += 1

    if added == 0:
        print("(!)  没有可用图片，未生成 .pptx")
        return None

    # 文件名用 plan title（去除非法字符）
    safe_title = re.sub(r"[^\w\u4e00-\u9fff\-]+", "_", title)[:60] or "deck"
    pptx_path = os.path.join(output_dir, f"{safe_title}.pptx")
    prs.save(pptx_path)
    print(f"  📑 PPTX generated: {pptx_path}  ({added} slides)")
    return pptx_path


# =============================================================================
# Commands (edit, rollback, ingest, list-sessions)
# =============================================================================

def cmd_list_sessions(args: argparse.Namespace) -> None:
    """List all generation sessions."""
    base_dir = args.output or str(CWD / OUTPUT_BASE_DIR)
    sessions = _find_sessions(base_dir)
    if not sessions:
        print("No sessions found.")
        return
    print(f"{'Session':<20} {'Slides':<8} Title")
    print("-" * 60)
    for s in sessions:
        print(f"{s['timestamp']:<20} {s['slide_count']:<8} {s['title']}")


def cmd_ingest_pptx(args: argparse.Namespace) -> None:
    """Ingest an external PPTX file and prepare it for editing.

    Renders the PPTX to PNGs, creates a session directory, and writes
    an initial metadata.json with placeholder slide_specs.  The Agent
    should then Read each page PNG and fill in the actual slide_specs.
    """
    pptx_path = args.ingest_pptx
    if not os.path.isfile(pptx_path):
        print(f"[X] File not found: {pptx_path}")
        sys.exit(1)

    # Render PPTX to PNGs
    sys.path.insert(0, str(SCRIPT_DIR))
    from render_template import render_pptx_to_pngs

    print(f"Ingesting: {pptx_path}")
    images_dir = render_pptx_to_pngs(pptx_path)
    print(f"Rendered to: {images_dir}")

    # Create session directory
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    if args.session:
        session_id = args.session
    else:
        stem = Path(pptx_path).stem
        session_id = f"{timestamp}_{stem}"
    session_dir = args.output or str(CWD / OUTPUT_BASE_DIR / session_id)
    session_images = os.path.join(session_dir, "images")
    os.makedirs(session_images, exist_ok=True)

    # Copy rendered PNGs as version 1 images
    import glob as _glob
    import shutil as _shutil

    png_files = sorted(_glob.glob(os.path.join(images_dir, "page-*.png")))
    slide_order: List[int] = []
    for i, png in enumerate(png_files, start=1):
        dest = os.path.join(session_images, f"slide-{i:02d}_v0001.png")
        _shutil.copy2(png, dest)
        # Also set as current
        current = os.path.join(session_images, f"slide-{i:02d}.png")
        _shutil.copy2(png, current)
        slide_order.append(i)

    # Build initial metadata with placeholder specs
    slides_meta: Dict[str, Any] = {}
    for i in slide_order:
        slides_meta[str(i)] = {
            "slide_number": i,
            "page_type": "content",
            "current_version": 1,
            "image_snapshot": f"images/slide-{i:02d}.png",
            "versions": [
                {
                    "version": 1,
                    "action": "ingest",
                    "spec": {
                        "layout": "(待 Agent 分析填充)",
                        "elements": {},
                    },
                    "prompt_file": "",
                    "image_snapshot": f"images/slide-{i:02d}_v0001.png",
                    "source_pptx": os.path.basename(pptx_path),
                }
            ],
        }

    metadata: Dict[str, Any] = {
        "version": 1,
        "title": Path(pptx_path).stem,
        "source_pptx": pptx_path,
        "slide_order": slide_order,
        "ingested_at": datetime.now().isoformat(),
        "slides": slides_meta,
    }
    _save_metadata(metadata, session_dir)

    print(f"Session created: {session_dir}")
    print(f"  {len(slide_order)} slides ingested")
    print(f"  metadata.json written with placeholder specs")
    print()
    print("Next: Agent should Read each page PNG and fill in slide_specs.")
    print("  For each slide, describe elements with type/content/position/style.")
    print("  Then use --edit to refine individual slides, or --plan to regenerate.")
    print(f"  Session timestamp: {session_id}")


def cmd_edit_slide(args: argparse.Namespace) -> None:
    """Edit a specific slide in an existing session.

    Reads the current slide_spec from metadata.json, applies element updates,
    constructs an edit prompt, and regenerates the slide with the original
    image as reference.
    """
    import shutil as _shutil

    session_dir = _resolve_session(args.session)
    metadata = _load_metadata(session_dir)
    images_dir = os.path.join(session_dir, "images")
    os.makedirs(images_dir, exist_ok=True)

    slide_key = str(args.edit)
    if slide_key not in metadata.get("slides", {}):
        print(f"[X] Slide {args.edit} not found in session {args.session}")
        sys.exit(1)

    slide_data = metadata["slides"][slide_key]
    current_version = slide_data["current_version"]
    page_type = slide_data.get("page_type", "content")

    # Get the current spec
    current_spec = _get_latest_slide_spec(metadata, args.edit)
    if not current_spec:
        print(f"[X] No spec found for slide {args.edit} (version {current_version})")
        print("    This session may not have slide_spec data. Consider regenerating.")
        sys.exit(1)

    # Apply element updates if provided
    element_updates: Dict[str, Any] = {}
    if args.element_updates:
        try:
            element_updates = json.loads(args.element_updates)
        except json.JSONDecodeError as e:
            print(f"[X] Invalid element-updates JSON: {e}")
            sys.exit(1)

    updated_spec = apply_spec_updates(current_spec, element_updates)

    # Construct edit prompt
    if args.edit_prompt:
        edit_prompt = args.edit_prompt
        if not element_updates:
            print("(!) Warning: --edit-prompt used without --element-updates.")
            print("    The slide_spec in metadata will NOT be updated.")
            print("    Future structured edits may reference stale content.")
    elif element_updates:
        edit_prompt = construct_edit_prompt(current_spec, element_updates)
    else:
        print("[X] No edit prompt or element updates provided")
        sys.exit(1)

    # Get reference image (current slide image)
    ref_path = os.path.join(session_dir, slide_data.get("image_snapshot", ""))
    if not os.path.isfile(ref_path):
        ref_path = os.path.join(images_dir, f"slide-{args.edit:02d}.png")
    if not os.path.isfile(ref_path):
        print(f"(!) Reference image not found at {ref_path}, generating without reference")
        ref_path = None

    # Backup current image and build new version
    new_version = current_version + 1
    current_img = os.path.join(images_dir, f"slide-{args.edit:02d}.png")
    versioned_img = os.path.join(images_dir, f"slide-{args.edit:02d}_v{current_version:04d}.png")
    if os.path.exists(current_img):
        _shutil.copy2(current_img, versioned_img)
    _stabilize_version_snapshots(slide_data, args.edit, images_dir)

    # Save edit prompt
    prompt_rel = f"images/slide-{args.edit:02d}_v{new_version:04d}.txt"
    prompt_abs = os.path.join(session_dir, prompt_rel)
    with open(prompt_abs, "w", encoding="utf-8") as f:
        f.write(edit_prompt)

    print(f"Editing slide {args.edit} (v{current_version} -> v{new_version})")
    if element_updates:
        for elem_id, changes in element_updates.items():
            for k, v in changes.items():
                old = current_spec.get("elements", {}).get(elem_id, {}).get(k, "?")
                print(f"  {elem_id}.{k}: {old} -> {v}")
    print(f"  Reference: {ref_path or '(none)'}")
    print(f"  Edit prompt: {edit_prompt[:200]}...")

    # Generate new image
    image_path = generate_slide(
        edit_prompt, args.edit, session_dir,
        reference_image_path=ref_path,
        backend=args.backend,
    )

    if not image_path:
        print("[X] Edit generation failed")
        sys.exit(1)

    # Also save versioned copy of the new image
    versioned_new = os.path.join(images_dir, f"slide-{args.edit:02d}_v{new_version:04d}.png")
    _shutil.copy2(current_img, versioned_new)

    # Update metadata
    _add_slide_version(
        slide_data=slide_data,
        new_version=new_version,
        spec=updated_spec,
        action="edit",
        image_snapshot=f"images/slide-{args.edit:02d}_v{new_version:04d}.png",
        prompt_file=prompt_rel,
        edit_instruction=args.edit_instruction or args.edit_prompt or "",
        reference_version=current_version,
    )

    _save_metadata(metadata, session_dir)

    # Regenerate PPTX
    slide_nums = _collect_slide_numbers(metadata)
    if not args.no_pptx:
        generate_pptx(session_dir, slide_nums, title=metadata.get("title", "Untitled"))

    print(f"Slide {args.edit} updated to v{new_version}")
    print(f"  Image: {image_path}")
    print(f"  Metadata saved")


def cmd_rollback_slide(args: argparse.Namespace) -> None:
    """Rollback a slide to a previous version.

    Reads the target version's spec, regenerates from it (optionally with
    the target version's image as reference), and saves as a new version.
    """
    import shutil as _shutil

    session_dir = _resolve_session(args.session)
    metadata = _load_metadata(session_dir)
    images_dir = os.path.join(session_dir, "images")

    slide_key = str(args.rollback)
    if slide_key not in metadata.get("slides", {}):
        print(f"[X] Slide {args.rollback} not found in session {args.session}")
        sys.exit(1)

    slide_data = metadata["slides"][slide_key]
    current_version = slide_data["current_version"]
    target_version = args.to_version

    target = _get_version_info(metadata, args.rollback, target_version)
    if not target:
        print(f"[X] Version {target_version} not found for slide {args.rollback}")
        available = sorted(v.get("version", 0) for v in slide_data.get("versions", []))
        print(f"    Available versions: {available}")
        sys.exit(1)

    target_spec = target.get("spec", {})
    target_image_rel = target.get("image_snapshot", "")

    # Construct generation prompt from target spec
    # Use style from metadata (stored during generation), fallback to --style arg
    style_template = ""
    stored_style = metadata.get("style", "")
    if stored_style:
        # Try the stored path; it may be relative to CWD or SCRIPT_DIR
        for candidate_path in (stored_style, str(SCRIPT_DIR / stored_style)):
            if os.path.isfile(candidate_path):
                style_template = load_style_template(candidate_path)
                break
    if not style_template and hasattr(args, "style") and args.style:
        style_template = load_style_template(args.style)

    if target_spec.get("elements"):
        prompt = generate_prompt_from_spec(
            style_template or "按以下元素描述生成幻灯片页",
            target_spec,
            slide_data.get("page_type", "content"),
            args.rollback,
            len(metadata.get("slides", {})),
        )
    else:
        # Fallback: use stored prompt if available
        prompt_path = os.path.join(session_dir, target.get("prompt_file", ""))
        if prompt_path and os.path.isfile(prompt_path):
            with open(prompt_path, "r", encoding="utf-8") as f:
                prompt = f.read()
        elif target_image_rel:
            # Ingested slide with placeholder spec: regenerate from reference image
            ref_abs = os.path.join(session_dir, target_image_rel)
            if os.path.isfile(ref_abs):
                print(f"  Spec has no elements; using reference image as visual guide")
                prompt = (
                    "请生成一张与参考图视觉风格完全一致的幻灯片页面。\n"
                    "保持相同的布局、配色、字体和装饰元素。\n"
                    "如果参考图上有文字，保留其内容和位置。"
                )
            else:
                print(f"[X] No elements in spec and no reference image found for version {target_version}")
                sys.exit(1)
        else:
            print(f"[X] No elements in spec and no prompt file or reference image for version {target_version}")
            print(f"    This often happens with ingested slides before the Agent fills in slide_specs.")
            print(f"    Have the Agent Read each page PNG and fill in the spec elements first.")
            sys.exit(1)

    # Get reference image from target version
    ref_path = os.path.join(session_dir, target_image_rel) if target_image_rel else None
    if ref_path and not os.path.isfile(ref_path):
        print(f"(!) Reference image not found, regenerating without reference")
        ref_path = None

    # Backup current image
    new_version = current_version + 1
    current_img = os.path.join(images_dir, f"slide-{args.rollback:02d}.png")
    versioned_img = os.path.join(images_dir, f"slide-{args.rollback:02d}_v{current_version:04d}.png")
    if os.path.exists(current_img):
        _shutil.copy2(current_img, versioned_img)
    _stabilize_version_snapshots(slide_data, args.rollback, images_dir)

    # Save prompt
    prompt_rel = f"images/slide-{args.rollback:02d}_v{new_version:04d}.txt"
    prompt_abs = os.path.join(session_dir, prompt_rel)
    with open(prompt_abs, "w", encoding="utf-8") as f:
        f.write(prompt)

    print(f"Rolling back slide {args.rollback}: v{current_version} -> v{new_version} (from v{target_version})")

    image_path = generate_slide(
        prompt, args.rollback, session_dir,
        reference_image_path=ref_path,
        backend=args.backend,
    )

    if not image_path:
        print("[X] Rollback generation failed")
        sys.exit(1)

    # Save versioned copy of new image
    versioned_new = os.path.join(images_dir, f"slide-{args.rollback:02d}_v{new_version:04d}.png")
    _shutil.copy2(current_img, versioned_new)

    # Update metadata
    _add_slide_version(
        slide_data=slide_data,
        new_version=new_version,
        spec=target_spec,
        action="rollback",
        image_snapshot=f"images/slide-{args.rollback:02d}_v{new_version:04d}.png",
        prompt_file=prompt_rel,
        edit_instruction=f"Rollback to version {target_version}",
        reference_version=target_version,
    )

    _save_metadata(metadata, session_dir)

    # Regenerate PPTX
    slide_nums = _collect_slide_numbers(metadata)
    if not args.no_pptx:
        generate_pptx(session_dir, slide_nums, title=metadata.get("title", "Untitled"))

    print(f"Slide {args.rollback} rolled back (v{new_version}, based on v{target_version})")


# =============================================================================
# Main Entry Point
# =============================================================================

def create_argument_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="PPT Generator - Generate PPT images using OpenAI gpt-image-2",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Example usage:
  # Generate
  python scripts/generate_ppt.py --plan slides_plan.json --style styles/gradient-glass.md
  python scripts/generate_ppt.py --plan slides_plan.json --style styles/clean-tech-blue.md --slides 1,3,5

  # Edit
  python scripts/generate_ppt.py --edit 3 --session 20240523_143052 --element-updates '{"subtitle":{"content":"新内容"}}'

  # Rollback
  python scripts/generate_ppt.py --rollback 3 --to-version 1 --session 20240523_143052

  # Ingest external PPTX
  python scripts/generate_ppt.py --ingest-pptx path/to/deck.pptx

  # List sessions
  python scripts/generate_ppt.py --list-sessions

Environment variables:
  OPENAI_BASE_URL:        Images API base URL (default: https://api.openai.com)
  OPENAI_API_KEY:         API key (required)
  GPT_IMAGE_MODEL_NAME:   Model name (default: gpt-image-2)
  GPT_IMAGE_QUALITY:      low / medium / high / auto (default: high)
""",
    )

    parser.add_argument("--plan", help="Path to slides plan JSON file")
    parser.add_argument("--style", help="Path to style template file (与模板输入二选一)")
    parser.add_argument("--output", help="Output directory path (default: outputs/TIMESTAMP)")
    parser.add_argument(
        "--template-pptx",
        help="用户的 .pptx 模板路径，启用「仿模板」模式",
    )
    parser.add_argument(
        "--template-images",
        help="模板每页 PNG 所在目录（强烈建议传，没有则只读 .pptx XML，不能跑 vision）",
    )
    parser.add_argument(
        "--template-profile",
        help="预先分析好的 TemplateProfile JSON。多模态 agent / 原生 Codex 可自己看图生成它，从而不需要 VISION_*。",
    )
    parser.add_argument(
        "--template-strict",
        action="store_true",
        help="高保真模式：把模板对应页作为 image reference 传给 gpt-image-2 出新图",
    )
    parser.add_argument(
        "--rebuild-template-cache",
        action="store_true",
        help="无视模板缓存重新跑 vision",
    )
    parser.add_argument(
        "--slides",
        help="Only generate specific slides, e.g. '1,3,5'",
    )
    parser.add_argument(
        "--concurrency",
        type=int,
        default=int(os.getenv("GPT_IMAGE_CONCURRENCY", "10")),
        help="并发请求数（默认 10，可用 GPT_IMAGE_CONCURRENCY 环境变量覆盖）",
    )
    parser.add_argument(
        "--no-pptx",
        action="store_true",
        help="不生成 .pptx 文件（默认会自动打包成 16:9 PPTX）",
    )
    parser.add_argument(
        "--backend",
        choices=["openai", "codex"],
        default=os.getenv("GPT_IMAGE_BACKEND", "openai"),
        help="图片生成后端：openai=直调 OpenAI API（需 OPENAI_API_KEY，默认）；"
             "codex=走本地 codex CLI（复用 codex 登录，无需在本 skill 配 key，但更慢）",
    )

    # Edit / rollback / ingest commands
    parser.add_argument(
        "--edit",
        type=int,
        metavar="SLIDE_NUMBER",
        help="Edit a specific slide (requires --session). Use with --edit-prompt or --element-updates.",
    )
    parser.add_argument(
        "--session",
        help="Session timestamp or path (for --edit / --rollback). e.g. '20240523_143052'",
    )
    parser.add_argument(
        "--edit-prompt",
        help="Full edit prompt to send to gpt-image-2 for the edited slide",
    )
    parser.add_argument(
        "--edit-instruction",
        help="Human-readable description of the edit (stored in metadata)",
    )
    parser.add_argument(
        "--element-updates",
        help="JSON mapping element_id -> {key: new_value} to update the slide_spec."
             " e.g. '{\"subtitle\": {\"content\": \"新副标题\"}}'",
    )
    parser.add_argument(
        "--rollback",
        type=int,
        metavar="SLIDE_NUMBER",
        help="Rollback a slide to a previous version (requires --session --to-version)",
    )
    parser.add_argument(
        "--to-version",
        type=int,
        help="Target version number for rollback",
    )
    parser.add_argument(
        "--ingest-pptx",
        metavar="PPTX_PATH",
        help="Ingest an external PPTX file, render to PNGs, and create a session for editing",
    )
    parser.add_argument(
        "--list-sessions",
        action="store_true",
        help="List all generation sessions",
    )

    return parser


def main() -> None:
    load_skill_env()

    parser = create_argument_parser()
    args = parser.parse_args()

    # ── 命令分发 ────────────────────────────────────────────
    if args.list_sessions:
        cmd_list_sessions(args)
        return

    if args.ingest_pptx:
        cmd_ingest_pptx(args)
        return

    if args.edit:
        if not args.session:
            parser.error("--edit requires --session")
        cmd_edit_slide(args)
        return

    if args.rollback:
        if not args.session:
            parser.error("--rollback requires --session")
        if not args.to_version:
            parser.error("--rollback requires --to-version")
        cmd_rollback_slide(args)
        return

    # ── 生成模式：必须提供 --plan ───────────────────────────
    if not args.plan:
        parser.error("必须传 --plan（生成模式），或使用 --edit / --rollback / --ingest-pptx / --list-sessions")

    # 校验：style 与 template source 至少有一个
    use_template = bool(args.template_pptx or args.template_images or args.template_profile)
    if not use_template and not args.style:
        parser.error("必须传 --style 或 --template-pptx / --template-images / --template-profile 至少其一")

    style_template = ""
    if args.style:
        style_path = args.style
        if not os.path.isabs(style_path):
            candidate = SCRIPT_DIR / style_path
            if candidate.exists():
                style_path = str(candidate)
        style_template = load_style_template(style_path)
    else:
        style_path = "(template-derived)"

    # 模板模式：跑 vision 拿 TemplateProfile（带缓存）
    template_profile: Optional[Dict[str, Any]] = None
    if use_template:
        sys.path.insert(0, str(SCRIPT_DIR))
        if args.template_profile:
            with open(args.template_profile, "r", encoding="utf-8") as f:
                template_profile = json.load(f)
            print(f"📦 使用预分析模板 profile: {args.template_profile}")
        else:
            # 只给了 .pptx 没给 PNG -> 自动渲染到 <cwd>/template_renders/<stem>/
            if args.template_pptx and not args.template_images:
                from render_template import render_pptx_to_pngs
                print(f"🖨️  --template-images 未指定，自动渲染 {args.template_pptx}")
                args.template_images = str(render_pptx_to_pngs(args.template_pptx))
            from template_analyzer import analyze_template
            try:
                template_profile = analyze_template(
                    pptx_path=args.template_pptx,
                    images_dir=args.template_images,
                    rebuild=args.rebuild_template_cache,
                )
            except ValueError as e:
                msg = str(e)
                if "VISION_BASE_URL" in msg or "VISION_API_KEY" in msg:
                    print("[X] 模板克隆需要先获得 TemplateProfile。")
                    print("    多模态 agent / 原生 Codex：请先看 template_renders/page-*.png，生成 profile JSON 后用 --template-profile 传入。")
                    print("    纯文本 agent（如 DeepSeek）：请配置 VISION_BASE_URL / VISION_API_KEY / VISION_MODEL_NAME。")
                raise
        if not template_profile.get("layouts"):
            print("(!)  模板分析未产出 layouts（缺 --template-images？），将回退到自由风格 prompt")
            template_profile = None
            if not args.style:
                print("[X] 模板输入不可用且未提供 --style，无法构造有效生成风格。")
                print("    请提供有效 --template-profile / VISION_*，或额外传 --style 作为 fallback。")
                sys.exit(1)
        elif args.template_strict:
            missing_refs = [
                lay.get("id", f"layout-{i + 1:02d}")
                for i, lay in enumerate(template_profile.get("layouts", []))
                if not lay.get("reference_image")
            ]
            if missing_refs:
                print("(!)  --template-strict 已启用，但部分 layout 没有 reference_image，相关页会退化为纯 prompt 仿作。")
                print(f"    缺 reference_image 的 layout: {', '.join(missing_refs[:8])}")

    with open(args.plan, "r", encoding="utf-8") as f:
        slides_plan = json.load(f)

    if args.output:
        output_dir = args.output
    else:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        # 默认输出到调用者当前工作目录，而不是 skill 安装目录
        output_dir = str(CWD / OUTPUT_BASE_DIR / timestamp)

    os.makedirs(os.path.join(output_dir, "images"), exist_ok=True)

    slides = slides_plan["slides"]
    total_slides = len(slides)

    if args.slides:
        target_nums = set(int(x.strip()) for x in args.slides.split(","))
        slides = [s for s in slides if s.get("slide_number") in target_nums]

    selected_slide_numbers = [s["slide_number"] for s in slides]

    print("=" * 60)
    print("PPT Generator (gpt-image-2) Started")
    print("=" * 60)
    print(f"Style: {style_path}")
    if template_profile:
        print(f"Template: {template_profile.get('source')} (hash={template_profile.get('source_hash')}, "
              f"{len(template_profile.get('layouts', []))} layouts)")
        print(f"Strict mode: {args.template_strict}")
    print(f"Slides: {len(slides)} / {total_slides}")
    print(f"Output: {output_dir}")
    print(f"Concurrency: {args.concurrency}")
    print(f"Backend: {args.backend}")
    print("=" * 60)
    print()

    prompts_data: Dict[str, Any] = {
        "metadata": {
            "title": slides_plan.get("title", "Untitled Presentation"),
            "total_slides": total_slides,
            "model": os.getenv("GPT_IMAGE_MODEL_NAME", "gpt-image-2"),
            "style": style_path,
            "template": template_profile.get("source") if template_profile else None,
            "template_strict": args.template_strict if template_profile else False,
            "generated_at": datetime.now().isoformat(),
        },
        "slides": [],
    }

    # ── metadata.json（slide_spec 版本历史）──────────────────
    # If output_dir reuses an existing session, merge rather than overwrite
    existing_meta_path = os.path.join(output_dir, METADATA_FILENAME)
    if os.path.isfile(existing_meta_path):
        metadata = _load_metadata(output_dir)
        existing_order = set(metadata.get("slide_order", []))
        new_order = [n for n in selected_slide_numbers if n not in existing_order]
        if new_order:
            metadata["slide_order"] = list(existing_order) + new_order
        metadata["generated_at"] = datetime.now().isoformat()
        print(f"Merging with existing metadata ({len(existing_order)} existing + {len(new_order)} new slides)")
    else:
        metadata = {
            "version": 1,
            "title": slides_plan.get("title", "Untitled Presentation"),
            "style": style_path,
            "model": os.getenv("GPT_IMAGE_MODEL_NAME", "gpt-image-2"),
            "generated_at": datetime.now().isoformat(),
            "slide_order": selected_slide_numbers,
            "slides": {},
        }

    if template_profile:
        from template_analyzer import (
            match_layout,
            coerce_fields,
            render_prompt_from_template,
            check_layout_reuse,
        )
        # layout 复用检测：在派发任务前打出建议，让用户决定是否中断
        reuse_warnings = check_layout_reuse(slides, template_profile)
        if reuse_warnings:
            print()
            print("=" * 60)
            print("📐 Layout 复用检测（建议尽量做到 1 page : 1 layout）")
            print("=" * 60)
            for w in reuse_warnings:
                print(w)
            print("=" * 60)
            print()

    # 收集所有待跑任务（跳过已存在的）
    pending_tasks = []
    for slide_info in slides:
        slide_number = slide_info["slide_number"]
        page_type = slide_info.get("page_type", "content")
        content_text = slide_info.get("content", "")
        slide_spec = slide_info.get("slide_spec")

        existing = os.path.join(output_dir, "images", f"slide-{slide_number:02d}.png")
        if os.path.exists(existing):
            print(f"Slide {slide_number}: already exists, skipping.")
            prompts_data["slides"].append({
                "slide_number": slide_number,
                "page_type": page_type,
                "content": content_text,
                "prompt": "(skipped - already exists)",
                "image_path": existing,
            })
            # Also add to metadata for skipped slides (only if not already present)
            if slide_spec and str(slide_number) not in metadata.get("slides", {}):
                metadata["slides"][str(slide_number)] = _init_slide_metadata(
                    slide_number, page_type, slide_spec,
                    prompt_file="",
                    image_path=f"images/slide-{slide_number:02d}.png",
                )
            continue

        reference_image = None
        if template_profile:
            layout = match_layout(slide_info, template_profile)
            if layout is None:
                # 模板未匹配 -> 回退到 style_template
                prompt = generate_prompt(
                    style_template, page_type, content_text, slide_number, total_slides,
                    slide_spec=slide_spec,
                )
                matched_layout_id = None
            else:
                fields = coerce_fields(slide_info, layout)
                prompt = render_prompt_from_template(
                    profile=template_profile,
                    layout=layout,
                    fields=fields,
                    language_rule=LANGUAGE_FONT_RULE.strip(),
                )
                matched_layout_id = layout.get("id")
                if args.template_strict:
                    reference_image = layout.get("reference_image")
        else:
            prompt = generate_prompt(
                style_template, page_type, content_text, slide_number, total_slides,
                slide_spec=slide_spec,
            )
            matched_layout_id = None

        pending_tasks.append({
            "slide_number": slide_number,
            "page_type": page_type,
            "content": content_text,
            "prompt": prompt,
            "reference_image": reference_image,
            "layout_id": matched_layout_id,
            "slide_spec": slide_spec,
        })

    if pending_tasks:
        from concurrent.futures import ThreadPoolExecutor, as_completed

        worker_count = max(1, min(args.concurrency, len(pending_tasks)))
        print(f"📦 派发 {len(pending_tasks)} 个任务到 {worker_count} 个并发 worker...\n")

        results: Dict[int, Optional[str]] = {}

        def _run(task):
            n = task["slide_number"]
            print(f">️  [slide {n}] start ({task['page_type']}{' / ref' if task.get('reference_image') else ''})")
            try:
                path = generate_slide(
                    task["prompt"], n, output_dir,
                    reference_image_path=task.get("reference_image"),
                    backend=args.backend,
                )
                if not path:
                    raise RuntimeError("generator returned empty image path")
                print(f"[OK] [slide {n}] done")
                return n, path
            except Exception as e:
                print(f"[X] [slide {n}] failed: {e}")
                return n, None

        with ThreadPoolExecutor(max_workers=worker_count) as executor:
            futures = [executor.submit(_run, t) for t in pending_tasks]
            for fut in as_completed(futures):
                n, path = fut.result()
                results[n] = path

        # 按原顺序写回 prompts_data，同时填充 metadata
        for task in pending_tasks:
            n = task["slide_number"]
            image_path = results.get(n)
            prompts_data["slides"].append({
                "slide_number": n,
                "page_type": task["page_type"],
                "content": task["content"],
                "layout_id": task.get("layout_id"),
                "reference_image": task.get("reference_image"),
                "prompt": task["prompt"],
                "image_path": image_path,
                "slide_spec": task.get("slide_spec"),
            })

            # metadata.json: save slide_spec + version history
            spec = task.get("slide_spec") or {}
            prompt_rel = f"images/slide-{n:02d}_v0001.txt"
            img_rel = f"images/slide-{n:02d}_v0001.png"

            if image_path:
                import shutil as _shutil
                # Save prompt text for version history
                prompt_abs = os.path.join(output_dir, prompt_rel)
                os.makedirs(os.path.dirname(prompt_abs), exist_ok=True)
                with open(prompt_abs, "w", encoding="utf-8") as pf:
                    pf.write(task["prompt"])
                versioned_initial = os.path.join(output_dir, img_rel)
                _shutil.copy2(image_path, versioned_initial)

                metadata["slides"][str(n)] = _init_slide_metadata(
                    n, task["page_type"], spec, prompt_rel, img_rel
                )
            # For failed slides, still record the spec but no image
            elif spec:
                metadata["slides"][str(n)] = _init_slide_metadata(
                    n, task["page_type"], spec, "", ""
                )

    # 按 slide_number 排序，保证 prompts.json 与播放顺序一致
    prompts_data["slides"].sort(key=lambda s: s["slide_number"])
    print()

    # Save both legacy prompts.json and new metadata.json
    save_prompts(output_dir, prompts_data)
    _save_metadata(metadata, output_dir)

    failed_slides = sorted(
        slide["slide_number"]
        for slide in prompts_data["slides"]
        if not slide.get("image_path")
    )
    if failed_slides:
        failed_str = ", ".join(str(n) for n in failed_slides)
        print(f"[X] 生成失败，未继续产出 PPTX。失败页：{failed_str}")
        print(f"    metadata.json 已保存，修复后可 --edit 继续。")
        sys.exit(1)

    pptx_path = None
    if not args.no_pptx:
        pptx_path = generate_pptx(
            output_dir,
            _collect_slide_numbers(metadata),
            title=slides_plan.get("title", "Untitled"),
        )

    print()
    print("=" * 60)
    print("Generation Complete!")
    print("=" * 60)
    print(f"Output directory: {output_dir}")
    print(f"Metadata:    {os.path.join(output_dir, 'metadata.json')}")
    if pptx_path:
        print(f"PPTX file:   {pptx_path}")
    print()


if __name__ == "__main__":
    main()
